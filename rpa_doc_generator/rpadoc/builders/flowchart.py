"""Render a process flow as native, editable PowerPoint flowchart shapes.

Nodes become real MSO flowchart auto-shapes (terminator, process, decision,
parallelogram, predefined-process) and edges become real elbow connectors that
are *glued* to the shapes via ``begin_connect`` / ``end_connect`` — so in
PowerPoint you can drag a box and the arrows follow.  Nothing here is an image.

Layout strategy
---------------
We do a layered top-to-bottom layout.  The "spine" (the longest forward chain
from start) is placed in a centre column; any node that hangs off a decision's
alternate branch is offset into a side column on the same row.  Connectors are
glued by the nearest face, so cross-row / back edges route as elbows.
"""

from __future__ import annotations

from typing import Dict, List

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from . import theme

SHAPE_FOR_TYPE = {
    "start": MSO_SHAPE.FLOWCHART_TERMINATOR,
    "end": MSO_SHAPE.FLOWCHART_TERMINATOR,
    "process": MSO_SHAPE.FLOWCHART_PROCESS,
    "decision": MSO_SHAPE.FLOWCHART_DECISION,
    "io": MSO_SHAPE.FLOWCHART_DATA,
    "subprocess": MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS,
}

# Connection-site indices for these auto-shapes: 0=top, 1=left, 2=bottom, 3=right.
TOP, LEFT, BOTTOM, RIGHT = 0, 1, 2, 3


def add_flowchart(deck, nodes: List[dict], edges: List[dict], *,
                  title: str = "Process Flow Diagram"):
    """Add one or more flowchart slides to ``deck`` for the given graph."""
    if not nodes:
        return

    layout = _layout(nodes, edges)
    # Split into slides so boxes never get unreadably small.
    rows_per_slide = 7
    max_row = max(p["row"] for p in layout.values())
    slices = [(r, min(r + rows_per_slide - 1, max_row))
              for r in range(0, max_row + 1, rows_per_slide)]

    for si, (r0, r1) in enumerate(slices):
        suffix = "" if len(slices) == 1 else f"  ({si + 1}/{len(slices)})"
        deck._slide = deck._new_slide()
        deck._chrome(title + suffix,
                     "Editable PowerPoint shapes & connectors — drag to refine",
                     section_label=deck.doc_type)
        _draw_slice(deck, nodes, edges, layout, r0, r1)
        _legend(deck)


# --------------------------------------------------------------------------- #
# layout
# --------------------------------------------------------------------------- #
def _layout(nodes: List[dict], edges: List[dict]) -> Dict[str, dict]:
    """Assign each node a (row, col). col 0 = centre spine, 1 = right side branch."""
    ids = [n["id"] for n in nodes]
    out_edges: Dict[str, List[str]] = {i: [] for i in ids}
    in_count: Dict[str, int] = {i: 0 for i in ids}
    for e in edges:
        if e["from"] in out_edges and e["to"] in in_count:
            out_edges[e["from"]].append(e["to"])
            in_count[e["to"]] += 1

    # Start node: explicit 'start' type, else first node with no inbound edges.
    start = next((n["id"] for n in nodes if n["type"] == "start"), None)
    if start is None:
        start = next((i for i in ids if in_count[i] == 0), ids[0])

    pos: Dict[str, dict] = {}
    row = 0
    visited = set()

    # Primary spine: walk forward following the first unvisited successor.
    cur = start
    while cur is not None and cur not in visited:
        visited.add(cur)
        pos[cur] = {"row": row, "col": 0}
        row += 1
        nxt = None
        for succ in out_edges[cur]:
            if succ not in visited:
                nxt = succ
                break
        cur = nxt

    # Remaining nodes (decision side-branches, orphans): place near a predecessor,
    # offset to the side column.
    for n in nodes:
        nid = n["id"]
        if nid in pos:
            continue
        pred_row = None
        for e in edges:
            if e["to"] == nid and e["from"] in pos:
                pred_row = pos[e["from"]]["row"]
                break
        r = (pred_row + 1) if pred_row is not None else row
        # find a free side slot at/after r
        while any(p["row"] == r and p["col"] == 1 for p in pos.values()):
            r += 1
        pos[nid] = {"row": r, "col": 1}
        row = max(row, r + 1)

    return pos


def _draw_slice(deck, nodes, edges, layout, r0, r1):
    node_by_id = {n["id"]: n for n in nodes}
    rows = r1 - r0 + 1

    area_top = Inches(1.55)
    area_bottom = Inches(6.75)
    area_h = area_bottom - area_top
    row_h = int(area_h / rows)

    box_w = Inches(3.0)
    box_h = min(Inches(0.85), Emu(int(row_h * 0.62)))
    centre_x = Inches(4.6)              # left edge of centre column box
    side_x = Inches(8.7)               # left edge of side column box

    shapes: Dict[str, object] = {}
    for nid, p in layout.items():
        if not (r0 <= p["row"] <= r1):
            continue
        node = node_by_id[nid]
        left = centre_x if p["col"] == 0 else side_x
        top = area_top + Emu(int((p["row"] - r0) * row_h)) + Emu(int((row_h - int(box_h)) / 2))
        shapes[nid] = _draw_node(deck, node, left, top, box_w, box_h)

    # Connectors (only when both endpoints are on this slide).
    for e in edges:
        a, b = e["from"], e["to"]
        if a in shapes and b in shapes:
            _connect(deck, shapes[a], shapes[b], layout[a], layout[b], e.get("label", ""))


def _draw_node(deck, node, left, top, width, height):
    ntype = node["type"]
    shp = deck._slide.shapes.add_shape(SHAPE_FOR_TYPE.get(ntype, MSO_SHAPE.FLOWCHART_PROCESS),
                                       left, top, width, height)
    color = theme.NODE_COLORS.get(ntype, theme.SECONDARY)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = theme.WHITE
    shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    label = node["label"]
    size = 11 if len(label) < 40 else (10 if len(label) < 70 else 8.5)
    r = p.add_run(); r.text = label
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = theme.FONT
    # gold/io needs dark text for contrast
    r.font.color.rgb = theme.DARK_TEXT if ntype == "io" else theme.WHITE
    return shp


def _connect(deck, a, b, pa, pb, label):
    # Choose faces to glue: forward (down) uses bottom->top; side branch uses
    # right->top/left; back-edges use right->right elbow.
    if pb["row"] > pa["row"]:
        if pa["col"] == pb["col"]:
            begin_site, end_site = BOTTOM, TOP
        elif pb["col"] > pa["col"]:
            begin_site, end_site = RIGHT, TOP
        else:
            begin_site, end_site = LEFT, TOP
    elif pb["row"] < pa["row"]:                  # loop-back / upward edge
        begin_site, end_site = RIGHT, RIGHT
    else:                                        # same row, side by side
        begin_site, end_site = RIGHT, LEFT

    conn = deck._slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW, a.left, a.top, b.left, b.top)
    try:
        conn.begin_connect(a, begin_site)
        conn.end_connect(b, end_site)
    except Exception:
        pass
    conn.line.color.rgb = theme.PRIMARY
    conn.line.width = Pt(1.5)
    _arrow_head(conn)

    if label:
        _edge_label(deck, a, b, label)


def _arrow_head(conn):
    """Add an arrowhead at the connector's end (python-pptx has no direct API)."""
    ln = conn.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = ln.makeelement(qn("a:tailEnd"), {})
        ln.append(tail)
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")


def _edge_label(deck, a, b, label):
    mid_left = int((a.left + b.left) / 2)
    mid_top = int((a.top + b.top) / 2)
    box = deck._slide.shapes.add_textbox(Emu(mid_left), Emu(mid_top) - Inches(0.12),
                                         Inches(0.8), Inches(0.28))
    box.fill.solid()
    box.fill.fore_color.rgb = theme.WHITE
    box.line.color.rgb = theme.LINE
    box.line.width = Pt(0.5)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = theme.SECONDARY
    r.font.name = theme.FONT


def _legend(deck):
    """A small legend strip explaining the shape colours, left column."""
    left = Inches(0.5)
    top = Inches(1.7)
    deck._add_text(left, top - Inches(0.32), Inches(3.4), Inches(0.3),
                   [("LEGEND", 11, True, theme.PRIMARY)])
    for i, (ntype, label) in enumerate(theme.NODE_LABELS.items()):
        y = top + Emu(int(Inches(0.5) * i))
        chip = deck._add_box(left, y, Inches(0.3), Inches(0.22),
                             fill=theme.NODE_COLORS[ntype])
        deck._add_text(left + Inches(0.4), y - Inches(0.03), Inches(3.0), Inches(0.32),
                       [(label, 9.5, False, theme.DARK_TEXT)],
                       anchor=MSO_ANCHOR.MIDDLE)
