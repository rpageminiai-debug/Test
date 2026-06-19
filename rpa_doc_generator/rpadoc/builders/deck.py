"""A thin, opinionated wrapper over python-pptx for building clean corporate decks.

Every slide produced here uses the blank layout and shapes we add by hand, so the
output is fully editable in PowerPoint (real text frames, real tables, real
auto-shapes — no images).
"""

from __future__ import annotations

from typing import List, Optional, Sequence

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import theme

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_TOP = Inches(1.5)            # below the header band
CONTENT_W = SLIDE_W - 2 * MARGIN


class Deck:
    """Builds a 16:9 presentation with a consistent header/footer system."""

    def __init__(self, doc_type: str, model: dict):
        self.model = model
        self.doc_type = doc_type            # "SDD" or "PDD"
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank = self.prs.slide_layouts[6]
        self._section_no = 0

    # ------------------------------------------------------------------ #
    # low-level helpers
    # ------------------------------------------------------------------ #
    def _new_slide(self):
        return self.prs.slides.add_slide(self._blank)

    @staticmethod
    def _fill(shape, color):
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    @staticmethod
    def _no_autosize(tf):
        # Stop python-pptx/PowerPoint from auto-growing the shape.
        tf.word_wrap = True

    def _add_box(self, left, top, width, height, fill=None, line=None, line_w=None):
        from pptx.enum.shapes import MSO_SHAPE
        shp = self._slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        if fill is not None:
            self._fill(shp, fill)
        else:
            shp.fill.background()
        if line is not None:
            shp.line.color.rgb = line
            shp.line.width = line_w or Pt(0.75)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _add_text(self, left, top, width, height, runs, *, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP):
        """runs: list of (text, size, bold, color, italic) tuples, one per paragraph."""
        box = self._slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
        for i, spec in enumerate(runs):
            text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
            italic = spec[4] if len(spec) > 4 else False
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = theme.FONT
            r.font.color.rgb = color
        return box

    # ------------------------------------------------------------------ #
    # header / footer chrome
    # ------------------------------------------------------------------ #
    def _chrome(self, title: str, subtitle: str = "", section_label: str = ""):
        # Top band
        self._add_box(0, 0, SLIDE_W, Inches(1.15), fill=theme.PRIMARY)
        self._add_box(0, Inches(1.15), SLIDE_W, Pt(3), fill=theme.SECONDARY)
        runs = [(title, 22, True, theme.WHITE)]
        if subtitle:
            runs.append((subtitle, 12, False, theme.LIGHT))
        self._add_text(MARGIN, Inches(0.18), SLIDE_W - 2 * MARGIN - Inches(2.0),
                       Inches(0.9), runs, anchor=MSO_ANCHOR.MIDDLE)
        if section_label:
            self._add_text(SLIDE_W - Inches(2.4), Inches(0.18), Inches(1.9),
                           Inches(0.8), [(section_label, 11, True, theme.LIGHT)],
                           align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        # Footer
        proc = self.model["project"]["process_name"]
        ver = self.model["project"]["version"]
        self._add_box(0, SLIDE_H - Inches(0.34), SLIDE_W, Pt(1), fill=theme.LINE)
        self._add_text(MARGIN, SLIDE_H - Inches(0.32), Inches(8),
                       Inches(0.28),
                       [(f"{self.doc_type}  |  {proc}  |  v{ver}", 8, False, theme.GREY_TEXT)],
                       anchor=MSO_ANCHOR.MIDDLE)
        self._add_text(SLIDE_W - Inches(3.6), SLIDE_H - Inches(0.32), Inches(3),
                       Inches(0.28),
                       [("Confidential", 8, False, theme.GREY_TEXT)],
                       align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # ------------------------------------------------------------------ #
    # public slide builders
    # ------------------------------------------------------------------ #
    def title_slide(self, doc_title: str):
        self._slide = self._new_slide()
        self._add_box(0, 0, SLIDE_W, SLIDE_H, fill=theme.PRIMARY)
        self._add_box(0, Inches(4.55), SLIDE_W, Pt(3), fill=theme.SECONDARY)
        p = self.model["project"]
        self._add_text(MARGIN, Inches(2.3), SLIDE_W - 2 * MARGIN, Inches(0.6),
                       [(self.doc_type + "  ·  " + doc_title, 16, True, theme.LIGHT)])
        self._add_text(MARGIN, Inches(2.9), SLIDE_W - 2 * MARGIN, Inches(1.5),
                       [(p["process_name"], 40, True, theme.WHITE)])
        meta = []
        if p["client"]:
            meta.append(("Client: " + p["client"], 14, False, theme.LIGHT))
        if p["department"]:
            meta.append(("Department: " + p["department"], 14, False, theme.LIGHT))
        meta.append((f"Version {p['version']}    |    {p['date']}", 14, False, theme.LIGHT))
        if p["author"]:
            meta.append(("Prepared by: " + p["author"], 14, False, theme.LIGHT))
        self._add_text(MARGIN, Inches(4.8), SLIDE_W - 2 * MARGIN, Inches(2), meta)
        self._add_text(MARGIN, SLIDE_H - Inches(0.7), SLIDE_W - 2 * MARGIN, Inches(0.4),
                       [("Robotic Process Automation  ·  Solution & Process Design", 11,
                         False, theme.SECONDARY)])

    def section_divider(self, title: str):
        self._section_no += 1
        self._slide = self._new_slide()
        self._add_box(0, 0, SLIDE_W, SLIDE_H, fill=theme.PRIMARY)
        self._add_box(MARGIN, Inches(3.0), Inches(0.9), Inches(1.5), fill=theme.SECONDARY)
        self._add_text(MARGIN, Inches(3.0), Inches(0.9), Inches(1.5),
                       [(f"{self._section_no:02d}", 40, True, theme.WHITE)],
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._add_text(Inches(1.8), Inches(3.0), SLIDE_W - Inches(2.4), Inches(1.5),
                       [(title, 32, True, theme.WHITE)], anchor=MSO_ANCHOR.MIDDLE)

    def toc_slide(self, sections: Sequence[str]):
        self._slide = self._new_slide()
        self._chrome("Table of Contents")
        col_left = MARGIN
        col_w = (CONTENT_W - Inches(0.4)) / 2
        per_col = (len(sections) + 1) // 2
        for i, name in enumerate(sections):
            col = 0 if i < per_col else 1
            row = i if col == 0 else i - per_col
            left = col_left + col * (col_w + Inches(0.4))
            top = CONTENT_TOP + Inches(0.1) + row * Inches(0.62)
            self._add_box(left, top + Inches(0.04), Inches(0.42), Inches(0.42),
                          fill=theme.LIGHT)
            self._add_text(left, top + Inches(0.04), Inches(0.42), Inches(0.42),
                           [(f"{i+1:02d}", 13, True, theme.PRIMARY)],
                           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._add_text(left + Inches(0.55), top, col_w - Inches(0.55), Inches(0.5),
                           [(name, 14, False, theme.DARK_TEXT)],
                           anchor=MSO_ANCHOR.MIDDLE)

    def bullets_slide(self, title: str, items: Sequence[str], *, subtitle: str = "",
                      intro: str = "", numbered: bool = False, two_col: bool = False,
                      empty_msg: str = "To be confirmed with the business."):
        items = list(items) or [empty_msg]
        # paginate if very long
        cap = 18 if two_col else 11
        pages = [items[i:i + cap] for i in range(0, len(items), cap)] or [[empty_msg]]
        for pi, page in enumerate(pages):
            self._slide = self._new_slide()
            ttl = title if len(pages) == 1 else f"{title}  ({pi+1}/{len(pages)})"
            self._chrome(ttl, subtitle, section_label=self.doc_type)
            top = CONTENT_TOP
            if intro and pi == 0:
                box = self._add_text(MARGIN, top, CONTENT_W, Inches(0.7),
                                     [(intro, 13, False, theme.GREY_TEXT, True)])
                top = CONTENT_TOP + Inches(0.75)
            self._render_bullets(page, top, numbered=numbered, two_col=two_col,
                                 start_index=pi * cap)

    def _render_bullets(self, items, top, *, numbered=False, two_col=False, start_index=0):
        if two_col:
            half = (len(items) + 1) // 2
            cols = [items[:half], items[half:]]
            col_w = (CONTENT_W - Inches(0.4)) / 2
            lefts = [MARGIN, MARGIN + col_w + Inches(0.4)]
            offset = 0
            for ci, col_items in enumerate(cols):
                self._bullet_column(col_items, lefts[ci], top, col_w, numbered,
                                    start_index + offset)
                offset += len(col_items)
        else:
            self._bullet_column(items, MARGIN, top, CONTENT_W, numbered, start_index)

    def _bullet_column(self, items, left, top, width, numbered, start_index):
        box = self._slide.shapes.add_textbox(left, top, width, SLIDE_H - top - Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        for i, text in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            p.line_spacing = 1.05
            marker = f"{start_index + i + 1}.  " if numbered else "▪  "
            r = p.add_run()
            r.text = marker
            r.font.size = Pt(13)
            r.font.bold = True
            r.font.color.rgb = theme.SECONDARY
            r.font.name = theme.FONT
            r2 = p.add_run()
            r2.text = str(text)
            r2.font.size = Pt(13)
            r2.font.color.rgb = theme.DARK_TEXT
            r2.font.name = theme.FONT

    def narrative_slide(self, title: str, blocks: List[tuple], *, subtitle: str = ""):
        """blocks: list of (heading, body) pairs."""
        self._slide = self._new_slide()
        self._chrome(title, subtitle, section_label=self.doc_type)
        top = CONTENT_TOP
        for heading, body in blocks:
            if not body:
                continue
            self._add_text(MARGIN, top, CONTENT_W, Inches(0.4),
                           [(heading, 15, True, theme.PRIMARY)])
            top = top + Inches(0.42)
            box = self._add_text(MARGIN, top, CONTENT_W, Inches(1.4),
                                 [(body, 13, False, theme.DARK_TEXT)])
            # estimate height ~ chars; keep simple, advance fixed-ish.
            lines = max(1, (len(body) // 110) + body.count("\n") + 1)
            top = top + Inches(0.26) * lines + Inches(0.25)
            if top > SLIDE_H - Inches(1.0):
                break

    def table_slide(self, title: str, headers: Sequence[str], rows: Sequence[Sequence[str]],
                    *, col_widths: Optional[Sequence[float]] = None, subtitle: str = "",
                    font_size: int = 10, rows_per_slide: int = 11,
                    empty_msg: str = "To be confirmed."):
        rows = [list(r) for r in rows]
        if not rows:
            rows = [[empty_msg] + [""] * (len(headers) - 1)]
        pages = [rows[i:i + rows_per_slide] for i in range(0, len(rows), rows_per_slide)]
        for pi, page in enumerate(pages):
            self._slide = self._new_slide()
            ttl = title if len(pages) == 1 else f"{title}  ({pi+1}/{len(pages)})"
            self._chrome(ttl, subtitle, section_label=self.doc_type)
            self._draw_table(headers, page, col_widths, font_size)

    def _draw_table(self, headers, rows, col_widths, font_size):
        n_rows = len(rows) + 1
        n_cols = len(headers)
        left, top = MARGIN, CONTENT_TOP
        width = CONTENT_W
        height = Inches(0.42) + Inches(0.42) * len(rows)
        height = min(height, SLIDE_H - top - Inches(0.5))
        gtable = self._slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = gtable.table
        table.first_row = False
        table.horz_banding = False
        # remove built-in style banding via explicit fills below
        if col_widths:
            total = sum(col_widths)
            for c, w in enumerate(col_widths):
                table.columns[c].width = int(width * (w / total))
        # header
        for c, htext in enumerate(headers):
            cell = table.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.PRIMARY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(4); cell.margin_right = Pt(4)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(htext)
            r.font.size = Pt(font_size + 0.5)
            r.font.bold = True
            r.font.color.rgb = theme.WHITE
            r.font.name = theme.FONT
        # body
        for ri, row in enumerate(rows, start=1):
            band = theme.LIGHTER if ri % 2 else theme.WHITE
            for c in range(n_cols):
                cell = table.cell(ri, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = band
                cell.vertical_anchor = MSO_ANCHOR.TOP
                cell.margin_left = Pt(4); cell.margin_right = Pt(4)
                cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = str(row[c]) if c < len(row) and row[c] is not None else ""
                r.font.size = Pt(font_size)
                r.font.color.rgb = theme.DARK_TEXT
                r.font.name = theme.FONT
                if c == 0:
                    r.font.bold = True

    def save(self, path: str):
        self.prs.save(path)
